from pptx import Presentation
from pptx.util import Inches


def extract_section(text, start_tag, end_tags=None):

    try:

        start = text.index(start_tag) + len(start_tag)

        end = len(text)

        if end_tags:

            for tag in end_tags:

                pos = text.find(tag, start)

                if pos != -1:

                    end = min(end, pos)

        return text[start:end].strip()

    except:

        return ""


def create_presentation(slides_content):

    prs = Presentation()

    slides = slides_content.split(
        "================================================"
    )

    for slide_text in slides:

        if not slide_text.strip():
            continue

        title = extract_section(
            slide_text,
            "TITLE:",
            ["NARRATION:"]
        )

        narration = extract_section(
            slide_text,
            "NARRATION:",
            ["VISUAL:"]
        )

        visual = extract_section(
            slide_text,
            "VISUAL:"
        )

        slide = prs.slides.add_slide(
            prs.slide_layouts[5]
        )

        # =====================
        # TITLE
        # =====================

        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.3),
            Inches(8),
            Inches(0.8)
        )

        title_box.text = title

        # =====================
        # NARRATION
        # =====================

        narration_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.2),
            Inches(6),
            Inches(2.5)
        )

        narration_box.text = narration

        # =====================
        # VISUAL BOX
        # =====================

        visual_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(4),
            Inches(8),
            Inches(2)
        )

        visual_box.text = (
            "VISUAL SUMMARY\n\n"
            + visual
        )

    file_name = (
        "AI_Requirement_Presentation_V2.pptx"
    )

    prs.save(
        file_name
    )

    return file_name
